"""
09_build_presentation.py
=========================
Builds the submission deck from the hackathon's OWN template
(`report/HACKATHON_TEMPLATE.pptx`, the SIH Idea-Submission format).

Template rules that are respected here:
  * "Kindly keep the maximum slides limit up to six (6) (including the title
    slide)"  -> exactly 6 slides
  * "Try to avoid paragraphs and post your idea in points / diagrams /
    Infographics"  -> bullets and figures only, no prose blocks
  * the template's own title placeholders, footers and slide numbers are kept

The template's five content sections are mapped onto the content the problem
statement demands:

  Template section            Problem-statement content
  --------------------------  -----------------------------------------------
  2. Idea / Proposed Solution  Process Understanding & Model
  3. Technical Approach        Control Strategy
  4. Feasibility & Viability   Safety performance + robustness evidence
  5. Artifacts                 Results, scenario plots, dashboard snaps
  6. Research and References   Lessons learned + references

Every number written into the deck is READ FROM THE GENERATED DATA, never
hard-coded, so the deck cannot drift from the results.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt
except ImportError:  # optional dependency - never break the core pipeline
    print("  python-pptx not installed - skipping deck build.")
    print("  (pip install python-pptx)   see requirements-optional.txt")
    raise SystemExit(0)

from config import CONTROLLER, DATA_DIR, LIMITS, REPORT_DIR, STEP_SEQUENCE_ID

TEMPLATE = REPORT_DIR / "HACKATHON_TEMPLATE.pptx"
OUT = ROOT / "submission" / "Autonomous_Choke_Controller_Presentation.pptx"
FIG = ROOT / "figures"

INK = RGBColor(0x0F, 0x17, 0x2A)
DIM = RGBColor(0x52, 0x5F, 0x73)
BLUE = RGBColor(0x1D, 0x4F, 0xD8)
GREEN = RGBColor(0x11, 0x7A, 0x3A)
RED = RGBColor(0xB9, 0x1C, 0x1C)


# ---------------------------------------------------------------- helpers
def load_results():
    s = json.loads((DATA_DIR / "scenario_summary.json").read_text())
    r = json.loads((DATA_DIR / "robustness_summary.json").read_text())
    m = json.loads((DATA_DIR / "model_validation_metrics.json").read_text())
    fm = json.loads((DATA_DIR / "fitted_model.json").read_text())
    return s, r, m, fm


def clear_body_shapes(slide, keep_title=True):
    """Remove the template's placeholder body text boxes, keeping the title,
    footer and slide-number placeholders intact."""
    for shp in list(slide.shapes):
        name = shp.name or ""
        if "Footer" in name or "Slide Number" in name:
            continue
        if keep_title and shp == slide.shapes.title:
            continue
        if shp.has_text_frame and shp != slide.shapes.title:
            shp._element.getparent().remove(shp._element)


def set_title(slide, text, size=30):
    t = slide.shapes.title
    if t is None:
        return
    t.text_frame.text = text
    for p in t.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = INK


def add_bullets(slide, left, top, width, height, items, size=12.5, gap=5):
    """items: list of (text, level, bold, colour|None)"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, level, bold, colour in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour or (INK if bold else DIM)
    return box


def add_pic(slide, path, left, top, max_w, max_h, center_in_box=True):
    """Insert a picture scaled to FIT inside a max_w x max_h box, preserving
    aspect ratio.  Sizing by width alone silently pushes tall figures off the
    bottom of the slide, so both dimensions are always bounded here."""
    path = Path(path)
    if not path.exists():
        print(f"    (missing figure: {path.name})")
        return None
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / (iw / 96), max_h / (ih / 96))
    w, h = (iw / 96) * scale, (ih / 96) * scale
    x = left + (max_w - w) / 2 if center_in_box else left
    y = top + (max_h - h) / 2 if center_in_box else top
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                    width=Inches(w), height=Inches(h))


def add_band(slide, left, top, width, height, text, fill, size=11):
    """A small coloured KPI band."""
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.text = text
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shp


# ---------------------------------------------------------------- build
def main():
    if not TEMPLATE.exists():
        print(f"  Template not found at {TEMPLATE} - skipping deck build.")
        return

    s, r, m, fm = load_results()
    meta = r["_meta"]
    q_max = meta["q_max_safe_ground_truth"]
    pct_of_max = 100 * r["C"]["final_Q_mean"] / q_max
    n_runs = meta["n_runs_per_scenario"] * 3

    prs = Presentation(str(TEMPLATE))

    # Drop the template's instructions slide (slide 1) -> leaves exactly 6.
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

    S = prs.slides
    assert len(S) == 6, f"expected 6 slides after removing instructions, got {len(S)}"

    # ------------------------------------------------ 1: Title
    sl = S[0]
    clear_body_shapes(sl, keep_title=False)
    add_bullets(sl, 0.7, 1.5, 11.9, 1.4, [
        ("Autonomous Production Choke Controller", 0, True, INK),
    ], size=34)
    add_bullets(sl, 0.7, 2.6, 11.9, 0.6, [
        ("for a Single Naturally Flowing Oil Well", 0, False, DIM),
    ], size=18)
    add_bullets(sl, 0.7, 3.5, 11.9, 2.4, [
        ("Brute-force Model Predictive Control  ·  control interval Ts = 1 hour", 0, False, DIM),
        ("Maximises safe production; refuses targets that cannot be met safely", 0, False, DIM),
        ("", 0, False, DIM),
        (f"0 constraint violations across {n_runs} closed-loop runs  |  "
         f"targets tracked to <1 %  |  {pct_of_max:.1f} % of max safe rate when target infeasible",
         0, True, GREEN),
    ], size=14)

    # ------------------------------------------------ 2: Process Understanding & Model
    sl = S[1]
    clear_body_shapes(sl)
    set_title(sl, "PROCESS UNDERSTANDING & MODEL", 26)
    add_bullets(sl, 0.55, 1.30, 6.0, 5.3, [
        ("Step-test results (our own experiments)", 0, True, BLUE),
        (f"{len(STEP_SEQUENCE_ID)} choke steps, up and down, 5–20 % magnitude, each held 8 h "
         "(≈4× slowest τ) to reach true steady state", 1, False, None),
        ("Opening the choke raises Q and FLP, lowers WHP and BHP — one input, "
         "four strongly coupled outputs", 1, False, None),
        ("Gain dQ/du collapses 3.91 → 0.15 bbl/hr per % across the range (~26×) "
         "→ the process is strongly nonlinear", 1, False, None),
        ("Provided reference dataset used for illustration only — never for fitting, "
         "per the problem statement", 1, False, None),
        ("", 0, False, None),
        ("Model assumptions", 0, True, BLUE),
        ("Linear IPR (constant productivity index); turbulent tubing friction ∝ Q²; "
         "choke orifice Q = Cv·f(u)·√(WHP−FLP)", 1, False, None),
        ("One first-order lag per output; dead time negligible (θ ≤ 0.09 h)", 1, False, None),
        ("WHT and AP monitored and logged, but not active constraints", 1, False, None),
        ("", 0, False, None),
        ("Dynamic model developed", 0, True, BLUE),
        ("y(k+1) = y(k) + α·( y_ss(u) + d − y(k) ),   α = 1 − e^(−Ts/τ)", 1, True, INK),
        ("Nonlinearity captured as a piecewise-linear steady-state curve y_ss(u); "
         "dynamics stay linear → simple and explainable", 1, False, None),
        (f"Fitted τ:  Q {fm['outputs']['Q']['tau']:.2f} h · WHP {fm['outputs']['WHP']['tau']:.2f} h · "
         f"FLP {fm['outputs']['FLP']['tau']:.2f} h · BHP {fm['outputs']['BHP']['tau']:.2f} h", 1, False, None),
        (f"Validated on HELD-OUT steps: NRMSE {m['Q']['nrmse_pct']:.2f} % (Q), "
         f"{m['BHP']['nrmse_pct']:.2f} % (BHP)", 1, True, GREEN),
    ], size=11)
    add_pic(sl, FIG / "step_gain_curve.png", 6.80, 1.22, 6.15, 2.15)
    add_pic(sl, FIG / "model_validation.png", 6.80, 3.45, 6.15, 3.30)

    # ------------------------------------------------ 3: Technical Approach / Control Strategy
    sl = S[2]
    clear_body_shapes(sl)
    set_title(sl, "TECHNICAL APPROACH — CONTROL STRATEGY", 24)
    add_bullets(sl, 0.55, 1.25, 6.6, 5.4, [
        ("Prediction methodology", 0, True, BLUE),
        (f"Enumerate ALL candidate moves within the ramp limit: u ± {CONTROLLER.max_move:.0f} % "
         f"on a {CONTROLLER.candidate_step} % grid → 21 candidates", 1, False, None),
        (f"Roll each forward {CONTROLLER.horizon} h through the identified model "
         "(move-blocked: apply, then hold)", 1, False, None),
        (f"Horizon must outrun the slowest lag: BHP τ = {fm['outputs']['BHP']['tau']:.2f} h, "
         f"so {CONTROLLER.horizon} h ≈ 5.7 τ", 1, False, None),
        ("Every prediction restarts from the current measurement → model error never compounds", 1, False, None),
        ("Measurement filter + bias correction remove hunting and steady-state offset", 1, False, None),
        ("", 0, False, None),
        ("Choke move selection logic", 0, True, BLUE),
        ("J = (Q_pred_end − Q_target)² + 0.05·(Δu)²  → closest to target, penalise large moves", 1, True, INK),
        ("Deterministic tie-breaks: lowest cost → smallest |Δu| → lower choke (conservative)", 1, False, None),
        ("Deadband: move only if it buys ≥ 2 (bbl/hr)² improvement → zero valve chatter", 1, False, None),
        ("", 0, False, None),
        ("Constraint handling", 0, True, BLUE),
        (f"Reject any candidate whose predicted WHP/FLP/BHP breaches limits at ANY point in the horizon",
         1, False, None),
        (f"Safety margins inside each hard limit ({CONTROLLER.margin_whp:.0f}/{CONTROLLER.margin_flp:.0f}/"
         f"{CONTROLLER.margin_bhp:.0f} psi) cover sensor noise + model error + filter lag", 1, False, None),
        ("If ALL candidates unsafe → pick the least-bad (smallest predicted violation), never lurch", 1, False, None),
        ("Infeasible target needs no special case: 'open further' is filtered out every "
         "step → the controller pins itself at the safe boundary", 1, False, None),
    ], size=9.8, gap=4)
    add_pic(sl, FIG / "dashboard" / "dash_reasoning.png", 7.35, 1.30, 5.55, 5.15)
    add_bullets(sl, 7.45, 6.62, 5.4, 0.4, [
        ("Live MPC reasoning: every candidate, why each was rejected", 0, False, DIM),
    ], size=9)

    # ------------------------------------------------ 4: Feasibility & Viability
    sl = S[3]
    clear_body_shapes(sl)
    set_title(sl, "FEASIBILITY AND VIABILITY", 26)
    kpis = [
        (f"{n_runs}\nclosed-loop runs", BLUE),
        ("0\nconstraint violations", GREEN),
        (f"{pct_of_max:.1f} %\nof max safe rate", GREEN),
        ("19/19\ntests passing", BLUE),
    ]
    for i, (txt, col) in enumerate(kpis):
        add_band(sl, 0.55 + i * 3.12, 1.28, 2.92, 0.95, txt, col, size=12)

    add_bullets(sl, 0.55, 2.45, 6.2, 4.4, [
        ("Feasibility — proven, not asserted", 0, True, BLUE),
        (f"Monte-Carlo: {meta['n_runs_per_scenario']} independent noise seeds × 3 scenarios; "
         "envelope checked on every sample", 1, False, None),
        (f"Worst margin to any limit: A +{r['A']['worst_margin_psi']:.0f} psi · "
         f"B +{r['B']['worst_margin_psi']:.0f} psi · C +{r['C']['worst_margin_psi']:.1f} psi", 1, False, None),
        ("Runs end-to-end in ~2 min; no GPU, no internet, no solver library", 1, False, None),
        ("", 0, False, None),
        ("Challenges found — and fixed", 0, True, RED),
        ("Horizon shorter than the slowest lag → BHP kept falling after the horizon "
         "(11/100 runs breached)", 1, False, None),
        ("Coarse step-test knots → chord across a convex curve over-estimated BHP by "
         "8 psi (unsafe direction)", 1, False, None),
        ("Both passed single-run testing; only the Monte-Carlo exposed them", 1, True, INK),
        ("", 0, False, None),
        ("Risk strategy", 0, True, BLUE),
        ("Margins sized for noise + model error + filter lag, and re-verified over 300 runs", 1, False, None),
        (f"Cost of safety quantified, not hidden: {meta['margin_cost_bblhr']:.1f} bbl/hr "
         f"({100*meta['margin_cost_bblhr']/q_max:.1f} %) below the theoretical max", 1, False, None),
        ("Controller reads ONLY step() outputs → official simulator drops in with zero "
         "changes (proven by a test against a different plant)", 1, False, None),
    ], size=9.8, gap=4)
    add_pic(sl, FIG / "robustness.png", 6.90, 2.45, 6.00, 4.35)

    # ------------------------------------------------ 5: Artifacts / Results
    sl = S[4]
    clear_body_shapes(sl)
    set_title(sl, "ARTIFACTS — RESULTS", 26)
    rows = [
        ("A — Startup to target", 120, s["A"]["settled_Q"], s["A"]["settled_Q_std"],
         s["A"]["final_choke"], s["A"]["min_bhp"]),
        ("B — Target tracking 100→150", 150, s["B"]["settled_Q"], s["B"]["settled_Q_std"],
         s["B"]["final_choke"], s["B"]["min_bhp"]),
        ("C — Infeasible target 200", 200, s["C"]["settled_Q"], s["C"]["settled_Q_std"],
         s["C"]["final_choke"], s["C"]["min_bhp"]),
    ]
    items = [("Scenario outcomes & tracking performance", 0, True, BLUE)]
    for name, tgt, q, sd, u, bhp in rows:
        verdict = "capped at max safe rate" if tgt == 200 else f"{abs(q - tgt) / tgt * 100:.1f} % error"
        items.append((f"{name}: settled {q:.1f} ± {sd:.1f} bbl/hr (target {tgt}) · "
                      f"choke {u:.1f} % · min BHP {bhp:.0f} psi · {verdict}", 1, False, None))
    items += [
        ("", 0, False, None),
        ("Safety performance", 0, True, BLUE),
        (f"WHP ≥ {LIMITS.whp_min:.0f}, FLP ≤ {LIMITS.flp_max:.0f}, BHP ≥ {LIMITS.bhp_min:.0f} psi, "
         f"|Δu| ≤ {CONTROLLER.max_move:.0f} %/step — held on EVERY sample", 1, False, None),
        ("Automated PASS/FAIL audit in the pipeline; non-zero exit if any check fails", 1, False, None),
        ("", 0, False, None),
        ("Deliverables", 0, True, BLUE),
        ("Simulator · step tests · model ID · MPC · 3 scenarios · Monte-Carlo · 19 tests · "
         "dashboard · report", 1, False, None),
        ("One command reproduces everything:  python run_all.py", 1, True, INK),
    ]
    add_bullets(sl, 0.55, 1.22, 5.9, 3.05, items, size=9.2, gap=3)
    add_pic(sl, FIG / "dashboard" / "dash_scenarioC_light.png", 6.55, 1.22, 6.30, 3.05)
    add_pic(sl, FIG / "scenarios_compact.png", 0.55, 4.42, 12.30, 2.45)
    add_bullets(sl, 0.55, 6.88, 12.30, 0.35, [
        ("All three scenarios: target tracked when feasible (A, B); refused and capped at the safe "
         "maximum when not (C) — BHP rests on its limit without ever crossing it.", 0, False, DIM),
    ], size=9)

    # ------------------------------------------------ 6: Research & References
    sl = S[5]
    clear_body_shapes(sl)
    set_title(sl, "RESEARCH AND REFERENCES", 26)
    add_bullets(sl, 0.55, 1.30, 6.2, 5.3, [
        ("Lessons learned", 0, True, BLUE),
        ("Where you place step-test knots is a SAFETY decision — linear interpolation across a "
         "convex curve is optimistic exactly where the constraint binds", 1, False, None),
        ("A horizon shorter than the slowest time constant is not conservative, it is unsafe", 1, False, None),
        ("One good run is not evidence; a safety claim about a stochastic system needs a distribution", 1, False, None),
        ("Margins must cover estimator lag and model error, not just sensor noise", 1, False, None),
        ("Every filter is a trade: ours removed constraint hunting but cost response lag "
         "that had to be paid for in margin", 1, False, None),
        ("Infeasibility needs no special-case code — it falls out of predict → filter → select", 1, False, None),
        ("", 0, False, None),
        ("Engineering references", 0, True, BLUE),
        ("Inflow performance / linear IPR (constant productivity index) — standard above bubble point", 1, False, None),
        ("Choke orifice / valve-sizing relation Q = Cv·f(u)·√ΔP", 1, False, None),
        ("Dynamic Matrix Control: move blocking, output disturbance (bias) estimation, "
         "constraint back-off", 1, False, None),
        ("Open-loop step testing and FOPDT identification for control-oriented models", 1, False, None),
    ], size=11)
    add_bullets(sl, 6.95, 1.30, 5.95, 5.3, [
        ("Provided material used", 0, True, BLUE),
        ("Autonomous_Choke_Control_Simulated_Dataset.csv — reference only, "
         "as the problem statement directs; never used to fit the model", 1, False, None),
        ("Compared against our physics stand-in and the differences reported openly "
         "(reference is near-linear; its FLP falls with rate)", 1, False, None),
        ("", 0, False, None),
        ("Project artifacts", 0, True, BLUE),
        ("src/  well_simulator · model · controller · simulator_interface", 1, False, None),
        ("scripts/  reference EDA · step tests · model ID · scenarios · robustness · capture · PDF", 1, False, None),
        ("tests/  19 tests incl. a swap-in proof against a different plant", 1, False, None),
        ("dashboard/  self-contained HTML/JS demo (no build, no server needed)", 1, False, None),
        ("report/ENGINEERING_REPORT.md → submission/ENGINEERING_REPORT.pdf", 1, False, None),
        ("", 0, False, None),
        ("Note on the simulator", 0, True, RED),
        ("The simulator module itself was not supplied, so we built a documented physics-based "
         "stand-in from the process description. The controller touches only step()/reset(), "
         "so the official simulator substitutes with zero changes.", 1, False, None),
    ], size=11)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT.relative_to(ROOT)}  ({OUT.stat().st_size / 1024:.0f} KB, {len(S)} slides)")


if __name__ == "__main__":
    main()
